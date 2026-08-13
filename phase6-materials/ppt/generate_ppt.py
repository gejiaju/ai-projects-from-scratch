# -*- coding: utf-8 -*-
"""
generate_ppt.py —— 生成项目总结 PPT（约 18 页）

这个文件是干什么的：
    用 python-pptx 把整个备考项目的核心成果自动生成一份 PPT，
    结构：背景 → 四个项目 → 踩坑经验 → 总结。每个项目都有"原理 + 结果"。
    生成后可以直接用 PowerPoint 打开编辑。

跑完能看到什么：
    当前目录生成 项目总结.pptx（约 18 页）

怎么跑：
    python3 generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 配色
DARK = RGBColor(0x1F, 0x4E, 0x8C)   # 深蓝（标题）
ACCENT = RGBColor(0x2C, 0xA0, 0x2C)  # 绿（强调）
GRAY = RGBColor(0x55, 0x55, 0x55)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
    if note:
        p = body.add_paragraph()
        p.text = note
        p.level = 1
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 封面
    add_title_slide(prs, "从零构建 AI 项目集",
                    "计算机视觉 · 目标检测 · 生成模型 · 大模型微调\n电子信息工程 → AI 方向")

    # 2. 目录
    add_bullet_slide(prs, "目录", [
        "1. 个人背景与技术栈",
        "2. 项目一：MLP vs CNN（深度学习基础）",
        "3. 项目二：YOLOv8 目标检测复现",
        "4. 项目三：注意力机制消融实验",
        "5. 项目四：DDPM 扩散模型",
        "6. 项目五：模型部署与大模型微调",
        "7. 踩坑经验与总结",
    ])

    # 3. 背景
    add_bullet_slide(prs, "个人背景与技术栈", [
        "本科：电子信息工程（信号处理、嵌入式系统）",
        "硬件：RTX 5060 (8GB) + WSL2 Ubuntu",
        "框架：PyTorch / ultralytics / transformers",
        "核心优势：EE 的系统思维 → 精度-速度-功耗权衡",
    ])

    # 4. 项目总览
    add_bullet_slide(prs, "四个项目总览", [
        "项目一：MLP vs CNN —— 手写对比，理解卷积本质",
        "项目二：YOLOv8 —— 完整复现 + 结构图",
        "项目三：注意力消融 —— 严格对照实验",
        "项目四：DDPM + 部署 + LoRA —— 三个进阶方向",
    ])

    # 5-6. 项目一
    add_bullet_slide(prs, "项目一：MLP vs CNN（原理）", [
        "两个模型用完全相同的训练流程和超参数，只换结构",
        "MLP：图片拍平成 784 个数 → 全连接层",
        "CNN：卷积层 + 池化层 + 全连接层",
    ])
    add_bullet_slide(prs, "项目一：结果（参数更少反而更准）", [
        "MLP：参数量 109,386，准确率 98.02%",
        "CNN：参数量 20,490，准确率 99.15%",
        "原因：参数共享 + 局部感受野",
    ])

    # 7-8. 项目二
    add_bullet_slide(prs, "项目二：YOLOv8 结构（三段式）", [
        "Backbone（CSPDarknet）：C2f 模块提特征",
        "Neck（PAN-FPN）：融合多尺度，输出 P3/P4/P5",
        "Head：anchor-free 解耦头，分类回归分开",
        "三个尺度分别检测小/中/大目标",
    ])
    add_bullet_slide(prs, "项目二：训练结果", [
        "COCO128 微调 30 轮：mAP50 = 80.3%",
        "推理速度 1.6ms/张（GPU）",
        "保存了训练曲线、混淆矩阵、PR 曲线",
    ])

    # 9-10. 项目三
    add_bullet_slide(prs, "项目三：注意力消融（方法）", [
        "SE：通道注意力（哪些通道重要）",
        "CBAM：通道 + 空间注意力（哪些位置重要）",
        "注入式替换保证三组只差注意力这一个变量",
    ])
    add_bullet_slide(prs, "项目三：消融结果", [
        "baseline：80.28%  |  +SE：80.24%  |  +CBAM：80.24%",
        "COCO128 上注意力无提升",
        "结论：加模块 ≠ 一定有效，效果取决于数据",
    ])

    # 11-12. 项目四
    add_bullet_slide(prs, "项目四：DDPM 扩散模型", [
        "前向：逐步加噪 x_t = √(ᾱ_t)x_0 + √(1-ᾱ_t)ε",
        "反向：训练 U-Net 预测噪声，逐步去噪",
        "从零手写 U-Net + 时间步 embedding",
    ])
    add_bullet_slide(prs, "项目四：可视化（噪声→图像）", [
        "MNIST 训练 20 轮，loss 0.092 → 0.040",
        "核心展示：噪声逐步变成手写数字的演变图",
        "这是 Stable Diffusion 的底层原理",
    ])

    # 13-14. 项目五
    add_bullet_slide(prs, "项目五：模型部署", [
        "ONNX 导出 + PyTorch/ONNX 推理速度对比",
        "踩坑：onnxruntime 默认 CPU 版，反而慢",
        "理解：部署加速需要选对推理后端（TensorRT）",
    ])
    add_bullet_slide(prs, "项目五：LoRA 大模型微调", [
        "冻结原模型，只训练旁路低秩矩阵 A·B",
        "参数量从 d×d 降到 2×d×r，省上千倍",
        "微调数据：电子信息考研面试问答",
    ])

    # 15. 踩坑经验
    add_bullet_slide(prs, "踩坑经验（独立解决问题的能力）", [
        "预训练权重下载失败，静默 fallback 从零训练 → 学会 val 验证",
        "ultralytics 版本更新重命名模型 → 用 summary 核对",
        "optimizer=auto 从零训练 lr 过小 → 定位 batch scaling",
    ])

    # 16. 总结
    add_bullet_slide(prs, "总结", [
        "完整闭环：理解原理 → 动手复现 → 严谨改进 → 部署落地",
        "工程习惯：不默认、先验证",
        "EE 优势：把精度-速度权衡迁移到 AI",
    ])

    # 17. 未来规划
    add_bullet_slide(prs, "未来研究方向", [
        "模型高效推理（TensorRT、量化、剪枝）",
        "多模态大模型与落地应用",
    ])

    # 18. 谢谢
    add_title_slide(prs, "谢谢", "请各位老师指正")

    out = "项目总结.pptx"
    prs.save(out)
    print(f"已生成 {out}，共 {len(prs.slides._sldIdLst)} 页")


if __name__ == "__main__":
    main()
